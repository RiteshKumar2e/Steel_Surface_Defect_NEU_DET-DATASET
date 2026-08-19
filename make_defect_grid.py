"""
NEU-DET: har defect type ka ek sample image, sabko ek hi grid figure me
bounding box ke saath -> PPT me paste karne ke liye ready PNG.

Usage:
    python make_defect_grid.py                 # default sample (pehli image)
    python make_defect_grid.py --random        # har class se random image
    python make_defect_grid.py --no-boxes      # bina bounding box ke
    python make_defect_grid.py --cols 3 --dpi 300
"""

import argparse
import random
import xml.etree.ElementTree as ET
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.patches as patches
import matplotlib.pyplot as plt
from PIL import Image

ROOT = Path(__file__).resolve().parent
IMAGES_DIR = ROOT / "NEU-DET" / "IMAGES"
ANNOT_DIR = ROOT / "NEU-DET" / "ANNOTATIONS"
OUT_DIR = ROOT / "OUTPUTS_ALL" / "defect_showcase"

# Folder name -> PPT me dikhne wala clean label
CLASS_LABELS = {
    "crazing": "Crazing",
    "inclusion": "Inclusion",
    "patches": "Patches",
    "pitted_surface": "Pitted Surface",
    "rolled-in_scale": "Rolled-in Scale",
    "scratches": "Scratches",
}

BOX_COLOR = "#ff2d55"


def load_boxes(stem):
    """Annotation XML se bounding boxes nikaalo. XML na mile to empty list."""
    xml_path = ANNOT_DIR / f"{stem}.xml"
    if not xml_path.exists():
        return []
    root = ET.parse(xml_path).getroot()
    boxes = []
    for obj in root.findall("object"):
        bb = obj.find("bndbox")
        if bb is None:
            continue
        boxes.append(
            (
                float(bb.findtext("xmin")),
                float(bb.findtext("ymin")),
                float(bb.findtext("xmax")),
                float(bb.findtext("ymax")),
            )
        )
    return boxes


def pick_sample(class_dir, use_random, seed):
    files = sorted(
        p for p in class_dir.iterdir() if p.suffix.lower() in {".jpg", ".jpeg", ".png", ".bmp"}
    )
    if not files:
        return None
    if use_random:
        return random.Random(seed).choice(files)
    return files[0]


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--random", action="store_true", help="har class se random image lo")
    ap.add_argument("--seed", type=int, default=42, help="random seed")
    ap.add_argument("--no-boxes", action="store_true", help="bounding box mat draw karo")
    ap.add_argument("--cols", type=int, default=3, help="grid columns")
    ap.add_argument("--dpi", type=int, default=300, help="output DPI (PPT ke liye 300 best)")
    ap.add_argument("--title", default="NEU-DET: Steel Surface Defect Types")
    ap.add_argument("--no-title", action="store_true", help="figure title hatao")
    ap.add_argument("--save-individual", action="store_true", help="har class ka alag PNG bhi")
    ap.add_argument("--pptx", action="store_true", help="ek .pptx slide bhi banao (python-pptx chahiye)")
    args = ap.parse_args()

    if not IMAGES_DIR.exists():
        raise SystemExit(f"IMAGES folder nahi mila: {IMAGES_DIR}")

    class_dirs = sorted(p for p in IMAGES_DIR.iterdir() if p.is_dir() and p.name in CLASS_LABELS)
    if not class_dirs:
        raise SystemExit(f"Koi defect class folder nahi mila {IMAGES_DIR} me")

    OUT_DIR.mkdir(parents=True, exist_ok=True)

    samples = []
    for class_dir in class_dirs:
        img_path = pick_sample(class_dir, args.random, args.seed)
        if img_path is None:
            print(f"[skip] {class_dir.name}: koi image nahi")
            continue
        samples.append((CLASS_LABELS[class_dir.name], img_path))

    cols = max(1, args.cols)
    rows = -(-len(samples) // cols)  # ceil

    fig, axes = plt.subplots(
        rows, cols, figsize=(cols * 3.2, rows * 3.6), layout="constrained"
    )
    axes = axes.ravel() if hasattr(axes, "ravel") else [axes]

    for ax, (label, img_path) in zip(axes, samples):
        img = Image.open(img_path).convert("RGB")
        ax.imshow(img)
        if not args.no_boxes:
            for xmin, ymin, xmax, ymax in load_boxes(img_path.stem):
                ax.add_patch(
                    patches.Rectangle(
                        (xmin, ymin),
                        xmax - xmin,
                        ymax - ymin,
                        linewidth=2,
                        edgecolor=BOX_COLOR,
                        facecolor="none",
                    )
                )
        ax.set_title(label, fontsize=13, fontweight="bold", pad=6)
        # filename image ke andar bottom-left me -> row spacing kabhi nahi bigadti
        ax.text(
            0.02,
            0.02,
            img_path.name,
            transform=ax.transAxes,
            fontsize=7,
            color="#222",
            va="bottom",
            ha="left",
            bbox=dict(facecolor="white", alpha=0.75, edgecolor="none", pad=1.5),
        )
        ax.set_xticks([])
        ax.set_yticks([])
        for spine in ax.spines.values():
            spine.set_edgecolor("#333")
            spine.set_linewidth(1.2)

    # bache hue khaali panels hatao
    for ax in axes[len(samples):]:
        ax.axis("off")

    if not args.no_title:
        fig.suptitle(args.title, fontsize=17, fontweight="bold")

    suffix = "_clean" if args.no_boxes else ""
    grid_path = OUT_DIR / f"all_defect_types_grid{suffix}.png"
    fig.savefig(grid_path, dpi=args.dpi, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    print(f"[saved] {grid_path}")

    if args.save_individual:
        for label, img_path in samples:
            f, a = plt.subplots(figsize=(3.2, 3.5))
            a.imshow(Image.open(img_path).convert("RGB"))
            if not args.no_boxes:
                for xmin, ymin, xmax, ymax in load_boxes(img_path.stem):
                    a.add_patch(
                        patches.Rectangle(
                            (xmin, ymin),
                            xmax - xmin,
                            ymax - ymin,
                            linewidth=2,
                            edgecolor=BOX_COLOR,
                            facecolor="none",
                        )
                    )
            a.set_title(label, fontsize=13, fontweight="bold")
            a.set_xticks([])
            a.set_yticks([])
            slug = label.lower().replace(" ", "_").replace("-", "_")
            out = OUT_DIR / f"{slug}{suffix}.png"
            f.savefig(out, dpi=args.dpi, bbox_inches="tight", facecolor="white")
            plt.close(f)
            print(f"[saved] {out}")

    if args.pptx:
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            print("[info] python-pptx nahi hai. Install: pip install python-pptx")
        else:
            prs = Presentation()
            prs.slide_width = Inches(13.333)  # 16:9
            prs.slide_height = Inches(7.5)
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = args.title
            pic_w = Inches(11)
            slide.shapes.add_picture(
                str(grid_path), (prs.slide_width - pic_w) // 2, Inches(1.5), width=pic_w
            )
            pptx_path = OUT_DIR / "defect_types_slide.pptx"
            prs.save(pptx_path)
            print(f"[saved] {pptx_path}")


if __name__ == "__main__":
    main()
